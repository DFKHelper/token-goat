from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "demo" / "fixtures" / "large-single-slide-review.pptx"


def main() -> None:
    width, height = 6000, 4000
    image = Image.effect_noise((width, height), 110).convert("RGB")
    image_buffer = BytesIO()
    image.save(image_buffer, format="JPEG", quality=95, subsampling=0)
    image_buffer.seek(0)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(image_buffer, 0, 0, width=presentation.slide_width, height=presentation.slide_height)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.6))
    title.text_frame.text = "Quarterly review: decision summary"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.55))
    body.text_frame.text = "Decision: approve the proposed rollout after the final controls review."
    body.text_frame.paragraphs[0].font.size = Pt(16)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT)
    print(f"{OUTPUT} ({OUTPUT.stat().st_size / 1024 / 1024:.1f} MB)")


if __name__ == "__main__":
    main()
